"""
PPTX Generator Module
Creates PowerPoint presentations with customizable themes and content
"""

import os
from typing import Dict, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from translations import get_text


def rgb(hexstr: str) -> RGBColor:
    """
    Convert hex color string to RGBColor.
    
    Args:
        hexstr: Hex color string (e.g., "#1F2D3D")
        
    Returns:
        RGBColor object
    """
    hexstr = hexstr.lstrip('#')
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))


def add_footer(slide, footer_text: str, show_numbers: bool, current: int, total: int, theme: Dict):
    """
    Add footer to a slide with optional slide numbers.
    
    Args:
        slide: PowerPoint slide object
        footer_text: Text to display in footer
        show_numbers: Whether to show slide numbers
        current: Current slide number
        total: Total number of slides
        theme: Theme configuration dictionary
    """
    ft = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3)).text_frame
    p = ft.paragraphs[0]
    p.text = f"{footer_text}  |  {current}/{total}" if show_numbers else footer_text
    p.font.size = Pt(10)
    p.font.color.rgb = rgb(theme["primary_color"])
    p.alignment = PP_ALIGN.LEFT


def add_title(slide, title: str, theme: Dict):
    """
    Add title to a slide.
    
    Args:
        slide: PowerPoint slide object
        title: Title text
        theme: Theme configuration dictionary
    """
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(1))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = rgb(theme["primary_color"])


def add_bullets(slide, bullets: List[str], theme: Dict, max_bullets: int = 6):
    """
    Add bullet points to a slide.
    
    Args:
        slide: PowerPoint slide object
        bullets: List of bullet point texts
        theme: Theme configuration dictionary
        max_bullets: Maximum number of bullets to display
    """
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(5.2), Inches(4.5))
    tf = tx.text_frame
    tf.clear()
    
    for idx, b in enumerate(bullets[:max_bullets]):
        para = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        para.text = b
        para.level = 0
        para.font.size = Pt(18)
        para.font.color.rgb = rgb(theme["primary_color"])


def add_image(slide, path: str, alt: str = ""):
    """
    Add image to a slide if it exists.
    
    Args:
        slide: PowerPoint slide object
        path: Path to image file
        alt: Alternative text for the image
    """
    if os.path.exists(path):
        try:
            slide.shapes.add_picture(path, Inches(6.1), Inches(1.7), width=Inches(3.2))
        except Exception as e:
            print(f"Warning: Could not add image {path}: {e}")


def add_logo(slide, logo_path: str, position: str = "top-right"):
    """
    Add logo to a slide.
    
    Args:
        slide: PowerPoint slide object
        logo_path: Path to logo file
        position: Position for the logo (top-right, top-left, etc.)
    """
    if os.path.exists(logo_path):
        try:
            if position == "top-right":
                slide.shapes.add_picture(logo_path, Inches(8.5), Inches(0.3), height=Inches(0.5))
            elif position == "top-left":
                slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.3), height=Inches(0.5))
        except Exception as e:
            print(f"Warning: Could not add logo {logo_path}: {e}")


def generate_pptx(spec: Dict, slides_content: List[Dict], out_path: str = "presentation.pptx", language: str = "es"):
    """
    Generate a complete PowerPoint presentation.
    
    Args:
        spec: Specification dictionary with theme, user info, etc.
        slides_content: List of slide content dictionaries
        out_path: Output file path
        language: Language code ("es" or "en")
    """
    prs = Presentation()
    total = spec["slides_count"]
    theme = spec["theme"]
    footer = spec["footer"]
    user = spec["user"]
    logo_path = user.get("logo_path", "")
    max_bullets = spec["length_target"].get("bullets_per_slide_max", 6)
    
    # Cover slide
    s0 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s0, user.get("brand_name", user.get("name", "Presentation")), theme)
    
    # Add subtitle if provided
    if spec.get("input", {}).get("description"):
        subtitle_box = s0.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
        p = subtitle_box.text_frame.paragraphs[0]
        desc = spec["input"]["description"]
        p.text = desc[:150] + "..." if len(desc) > 150 else desc
        p.font.size = Pt(16)
        p.font.color.rgb = rgb(theme.get("accent_color", theme["primary_color"]))
    
    if logo_path:
        add_logo(s0, logo_path)
    
    add_footer(s0, footer["text"], footer.get("show_slide_numbers", True), 1, total, theme)
    
    # Content slides
    for i, sc in enumerate(slides_content, start=2):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(s, sc["title"], theme)
        add_bullets(s, sc.get("bullets", []), theme, max_bullets)
        
        if sc.get("image_path"):
            add_image(s, sc["image_path"], sc.get("alt", ""))
        
        if logo_path:
            add_logo(s, logo_path)
        
        add_footer(s, footer["text"], footer.get("show_slide_numbers", True), i, total, theme)
    
    # Conclusion slide
    last = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(last, get_text("conclusions", language), theme)
    
    conclusion_bullets = [
        get_text("summary_learnings", language),
        get_text("immediate_actions", language),
        get_text("responsible_deadlines", language)
    ]
    add_bullets(last, conclusion_bullets, theme, max_bullets)
    
    if logo_path:
        add_logo(last, logo_path)
    
    add_footer(last, footer["text"], footer.get("show_slide_numbers", True), total, total, theme)
    
    # Save presentation
    prs.save(out_path)
    print(f"{get_text('file_saved', language)}: {out_path}")
